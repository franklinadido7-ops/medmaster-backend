"""
Extraction de texte depuis différents formats de documents + découpage en chunks.
"""
import re
from typing import List

import tiktoken
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from app.config import settings

_encoder = tiktoken.get_encoding("cl100k_base")


def extract_text(file_path: str, file_type: str) -> str:
    """Extrait le texte brut d'un fichier selon son type."""
    file_type = file_type.lower()

    if file_type == "pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if file_type in ("docx", "doc"):
        doc = DocxDocument(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if file_type in ("pptx", "ppt"):
        prs = Presentation(file_path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text)
        return "\n".join(chunks)

    if file_type in ("jpg", "jpeg", "png"):
        # Les images ne sont pas indexées dans le RAG textuel : elles sont
        # envoyées directement à Claude (vision) lors de la génération.
        return ""

    raise ValueError(f"Type de fichier non supporté pour l'extraction RAG : {file_type}")


def _token_len(text: str) -> int:
    return len(_encoder.encode(text))


def chunk_text(
    text: str,
    chunk_size: int = None,
    overlap: int = None,
) -> List[str]:
    """
    Découpe un texte en chunks d'environ `chunk_size` tokens, avec chevauchement
    `overlap`, en respectant si possible les frontières de paragraphes/phrases.
    """
    chunk_size = chunk_size or settings.CHUNK_SIZE
    overlap = overlap or settings.CHUNK_OVERLAP

    # Découpage initial par paragraphes
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    chunks: List[str] = []
    current = ""
    current_tokens = 0

    for para in paragraphs:
        para_tokens = _token_len(para)

        if current_tokens + para_tokens <= chunk_size:
            current = f"{current}\n\n{para}" if current else para
            current_tokens += para_tokens
        else:
            if current:
                chunks.append(current)
            # Si un seul paragraphe dépasse la taille max, le découper par phrases
            if para_tokens > chunk_size:
                sentences = re.split(r"(?<=[.!?])\s+", para)
                current, current_tokens = "", 0
                for sent in sentences:
                    st = _token_len(sent)
                    if current_tokens + st > chunk_size and current:
                        chunks.append(current)
                        current, current_tokens = "", 0
                    current = f"{current} {sent}".strip()
                    current_tokens += st
            else:
                current, current_tokens = para, para_tokens

    if current:
        chunks.append(current)

    # Application du chevauchement (overlap) entre chunks consécutifs
    if overlap > 0 and len(chunks) > 1:
        overlapped = [chunks[0]]
        for i in range(1, len(chunks)):
            prev_words = chunks[i - 1].split()
            overlap_words = prev_words[-overlap:] if len(prev_words) > overlap else prev_words
            overlapped.append(" ".join(overlap_words) + "\n\n" + chunks[i])
        chunks = overlapped

    return chunks
